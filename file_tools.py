"""
File Task Tools — 파일 경로 기반 요약·번역·분석 도구

지원 도구:
  load_file_content  : 파일 내용을 텍스트로 반환 (에이전트가 직접 열람)
  summarize_file     : 파일 내용 요약 (긴 문서 → 청크별 요약 후 통합)
  translate_file     : 파일 내용 번역 (긴 문서 → 청크별 번역 후 합산)
  extract_key_points : 핵심 포인트 추출
  analyze_file       : 커스텀 태스크 수행 (계약서 검토, 오류 찾기 등)

지원 파일 형식: PDF, CSV, 이미지, 텍스트
"""

import os
from pathlib import Path
from typing import Literal

from langchain.chat_models import init_chat_model
from langchain_core.tools import tool

from document_loader import load_document

# ── 쓰기 가능한 기본 출력 디렉터리 ─────────────────────────────────────
def _default_output_dir() -> Path:
    if env := os.environ.get("AGENT_DATA_DIR", "").strip():
        p = Path(env) / "reports"
    else:
        script_dir = Path(__file__).resolve().parent
        try:
            (script_dir / ".write_test").touch()
            (script_dir / ".write_test").unlink()
            p = script_dir / "reports"
        except OSError:
            p = Path("/tmp/deep_agent/reports")
    p.mkdir(parents=True, exist_ok=True)
    return p

_DEFAULT_OUTPUT_DIR = _default_output_dir()

# ── 처리 기준 상수 ──────────────────────────────────────────────────
MAX_DIRECT_CHARS = 12_000   # 이 이하면 청킹 없이 한 번에 처리
CHUNK_SIZE = 8_000          # 청킹 시 한 청크의 최대 글자 수


# ══════════════════════════════════════════════════════════════════
#  내부 헬퍼
# ══════════════════════════════════════════════════════════════════

def _load_text(file_path: str) -> str:
    """파일을 로드해 전체 텍스트를 반환. 오류 시 [Error] 접두 문자열 반환."""
    path = Path(file_path)
    if not path.exists():
        return f"[Error] 파일을 찾을 수 없습니다: {file_path}"

    result = load_document(path)
    if result.error:
        return f"[Error] 파일 로드 실패: {result.error}"

    text = result.to_text().strip()
    if not text:
        return f"[Error] 파일에서 텍스트를 추출할 수 없습니다: {file_path}"

    return text


def _chunk_text(text: str, chunk_size: int = CHUNK_SIZE) -> list[str]:
    """텍스트를 chunk_size 단위로 분할. 가능하면 줄바꿈 경계에서 자른다."""
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        if end < len(text):
            boundary = text.rfind("\n", start, end)
            if boundary > start:
                end = boundary
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end
    return chunks


def _call_llm(prompt: str) -> str:
    """gpt-4o를 호출해 응답 텍스트를 반환한다."""
    model = init_chat_model("openai:gpt-4o", temperature=0)
    return model.invoke(prompt).content


def _map_reduce(
    chunks: list[str],
    map_prompt_fn,
    reduce_prompt_fn,
) -> str:
    """
    청크별로 map_prompt_fn을 적용해 중간 결과를 만들고,
    reduce_prompt_fn으로 최종 결과를 합산한다.
    """
    partial_results: list[str] = []
    total = len(chunks)
    for i, chunk in enumerate(chunks, 1):
        print(f"[FileTools] 청크 {i}/{total} 처리 중...")
        partial_results.append(_call_llm(map_prompt_fn(chunk, i, total)))

    combined = "\n\n".join(
        f"[파트 {i}]\n{r}" for i, r in enumerate(partial_results, 1)
    )
    return _call_llm(reduce_prompt_fn(combined))


# ══════════════════════════════════════════════════════════════════
#  LangChain 도구 정의
# ══════════════════════════════════════════════════════════════════

@tool
def load_file_content(file_path: str) -> str:
    """
    파일 경로를 받아 텍스트 내용을 그대로 반환한다.
    PDF, CSV, 이미지, 텍스트 파일을 지원한다.
    에이전트가 파일 내용을 직접 확인하거나 인용할 때 사용.
    매우 긴 파일은 앞부분만 반환하며 summarize_file 사용을 권장한다.
    """
    text = _load_text(file_path)
    if text.startswith("[Error]"):
        return text

    if len(text) > 20_000:
        preview = text[:20_000]
        return (
            f"{preview}\n\n"
            f"[파일이 깁니다. 전체 {len(text):,}자 중 앞 20,000자를 반환했습니다. "
            f"전체 요약이 필요하면 summarize_file을 사용하세요.]"
        )
    return text


@tool
def summarize_file(
    file_path: str,
    style: Literal["bullet_points", "paragraph", "structured"] = "structured",
    language: str = "Korean",
) -> str:
    """
    파일 내용을 요약한다.
    긴 문서는 청크별로 요약한 뒤 통합해 최종 요약을 생성한다.

    Args:
        file_path : 요약할 파일 경로
        style     : bullet_points(글머리기호) | paragraph(단락) | structured(구조화 헤더)
        language  : 출력 언어 (예: Korean, English, Japanese)
    """
    text = _load_text(file_path)
    if text.startswith("[Error]"):
        return text

    style_guide = {
        "bullet_points": "글머리기호(•)로 핵심 내용을 나열하라.",
        "paragraph":     "자연스러운 단락 형식으로 서술하라.",
        "structured":    "## 개요 / ## 핵심 내용 / ## 결론 헤더를 사용해 구조화하라.",
    }.get(style, "구조화된 형식으로 요약하라.")

    if len(text) <= MAX_DIRECT_CHARS:
        prompt = (
            f"다음 문서를 {language}로 요약하라. {style_guide}\n\n"
            f"[문서]\n{text}"
        )
        return _call_llm(prompt)

    print(f"[FileTools] 긴 문서({len(text):,}자) — 청크 분할 요약 시작")
    chunks = _chunk_text(text)

    def map_fn(chunk, i, total):
        return (
            f"다음은 긴 문서의 일부({i}/{total})이다. "
            f"이 부분의 핵심을 {language}로 간략히 요약하라.\n\n{chunk}"
        )

    def reduce_fn(combined):
        return (
            f"아래는 긴 문서를 파트별로 요약한 결과다. "
            f"이를 통합해 하나의 최종 요약을 {language}로 작성하라. {style_guide}\n\n"
            f"{combined}"
        )

    return _map_reduce(chunks, map_fn, reduce_fn)


@tool
def translate_file(
    file_path: str,
    target_language: str = "Korean",
    source_language: str = "auto",
) -> str:
    """
    파일 내용을 지정한 언어로 번역한다.
    긴 문서는 청크별로 번역 후 순서대로 합친다.
    원문의 제목·단락·표 등 형식을 최대한 유지한다.

    Args:
        file_path       : 번역할 파일 경로
        target_language : 번역 목표 언어 (예: Korean, English, Japanese)
        source_language : 원본 언어 (auto면 자동 감지)
    """
    text = _load_text(file_path)
    if text.startswith("[Error]"):
        return text

    src_hint = "" if source_language == "auto" else f"원본 언어: {source_language}. "
    format_note = "원문의 제목, 단락, 표, 목록 등 형식을 유지하라."

    if len(text) <= MAX_DIRECT_CHARS:
        prompt = (
            f"{src_hint}다음 텍스트를 {target_language}로 번역하라. {format_note}\n\n"
            f"{text}"
        )
        return _call_llm(prompt)

    print(f"[FileTools] 긴 문서({len(text):,}자) — 청크 분할 번역 시작")
    chunks = _chunk_text(text)
    translated: list[str] = []
    total = len(chunks)
    for i, chunk in enumerate(chunks, 1):
        print(f"[FileTools] 청크 {i}/{total} 번역 중...")
        prompt = (
            f"{src_hint}다음 텍스트를 {target_language}로 번역하라. {format_note}\n\n"
            f"{chunk}"
        )
        translated.append(_call_llm(prompt))

    return "\n\n".join(translated)


@tool
def extract_key_points(
    file_path: str,
    max_points: int = 10,
    language: str = "Korean",
) -> str:
    """
    파일에서 핵심 포인트와 중요 정보를 번호 목록으로 추출한다.
    보고서, 계약서, 논문 등에서 핵심을 빠르게 파악할 때 사용.

    Args:
        file_path  : 분석할 파일 경로
        max_points : 추출할 최대 핵심 포인트 수
        language   : 출력 언어
    """
    text = _load_text(file_path)
    if text.startswith("[Error]"):
        return text

    if len(text) > MAX_DIRECT_CHARS:
        text = text[:MAX_DIRECT_CHARS] + "\n...(이하 생략, 앞부분 기준으로 추출)"

    prompt = (
        f"다음 문서에서 가장 중요한 핵심 포인트를 최대 {max_points}개 추출하라. "
        f"{language}로 작성하며, 각 항목에 번호를 붙여 명확히 서술하라.\n\n"
        f"[문서]\n{text}"
    )
    return _call_llm(prompt)


@tool
def analyze_file(file_path: str, task: str) -> str:
    """
    파일에 대해 사용자가 지정한 커스텀 작업을 수행한다.
    요약·번역 외의 분석(계약서 위험 조항 검토, 오류 찾기, 특정 질문 답변 등)에 사용.

    Args:
        file_path : 분석할 파일 경로
        task      : 수행할 작업 설명 (예: "이 계약서에서 갑에게 불리한 조항을 찾아줘")
    """
    text = _load_text(file_path)
    if text.startswith("[Error]"):
        return text

    if len(text) > MAX_DIRECT_CHARS:
        text = text[:MAX_DIRECT_CHARS] + "\n...(이하 생략)"

    prompt = (
        f"다음 문서에 대해 아래 작업을 수행하라.\n\n"
        f"[작업]\n{task}\n\n"
        f"[문서]\n{text}"
    )
    return _call_llm(prompt)


# ══════════════════════════════════════════════════════════════════
#  보고서 저장 — 형식별 내부 핸들러
# ══════════════════════════════════════════════════════════════════

def _parse_md_lines(content: str) -> list[dict]:
    """
    마크다운 텍스트를 줄 단위로 파싱해 딕셔너리 목록 반환.
    type: 'heading' | 'bullet' | 'text' | 'blank' | 'table_row' | 'table_sep'

    table_row  → cells 필드에 셀 값 목록 포함
    table_sep  → |---|---| 구분선 (헤더/데이터 경계 표시용)
    """
    result = []
    for line in content.splitlines():
        stripped = line.rstrip()
        if stripped.startswith("### "):
            result.append({"type": "heading", "level": 3, "text": stripped[4:], "cells": []})
        elif stripped.startswith("## "):
            result.append({"type": "heading", "level": 2, "text": stripped[3:], "cells": []})
        elif stripped.startswith("# "):
            result.append({"type": "heading", "level": 1, "text": stripped[2:], "cells": []})
        elif stripped.startswith("|"):
            inner = stripped.strip("|")
            # 구분선 감지: |---|---| 또는 |:---:|:---:| 패턴
            if "-" in inner and all(c in "-:| " for c in inner):
                result.append({"type": "table_sep", "level": 0, "text": "", "cells": []})
            else:
                cells = [c.strip() for c in inner.split("|")]
                result.append({"type": "table_row", "level": 0, "text": stripped, "cells": cells})
        elif stripped.startswith(("- ", "• ", "* ")):
            result.append({"type": "bullet", "level": 1, "text": stripped[2:], "cells": []})
        elif stripped == "":
            result.append({"type": "blank", "level": 0, "text": "", "cells": []})
        else:
            result.append({"type": "text", "level": 0, "text": stripped, "cells": []})
    return result


def _save_txt(content: str, file_path: Path) -> None:
    """txt / md: 텍스트 그대로 저장."""
    file_path.write_text(content, encoding="utf-8")


def _save_docx(content: str, file_path: Path) -> None:
    """
    docx: 마크다운 구조를 Word 스타일로 변환.
    # → Heading 1, ## → Heading 2, ### → Heading 3
    - / • → List Bullet
    | table | → Word 실제 테이블 (첫 행 헤더 스타일)
    나머지 → Normal
    """
    from docx import Document
    from docx.shared import Pt
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()
    doc.styles["Normal"].font.name = "맑은 고딕"
    doc.styles["Normal"].font.size = Pt(11)

    def _set_cell(cell, text: str, bold: bool = False, bg_hex: str | None = None) -> None:
        """셀 텍스트 설정 + 옵션 헤더 스타일 적용."""
        para = cell.paragraphs[0]
        para.clear()
        run = para.add_run(text)
        if bold:
            run.bold = True
        if bg_hex:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:val"), "clear")
            shd.set(qn("w:color"), "auto")
            shd.set(qn("w:fill"), bg_hex)
            tcPr.append(shd)

    items = _parse_md_lines(content)
    i = 0
    while i < len(items):
        item = items[i]

        if item["type"] == "heading":
            doc.add_heading(item["text"], level=item["level"])
            i += 1

        elif item["type"] == "bullet":
            doc.add_paragraph(item["text"], style="List Bullet")
            i += 1

        elif item["type"] == "blank":
            doc.add_paragraph("")
            i += 1

        elif item["type"] == "table_row":
            # 연속된 table_row / table_sep 블록을 한꺼번에 수집
            table_rows: list[list[str]] = []
            while i < len(items) and items[i]["type"] in ("table_row", "table_sep"):
                if items[i]["type"] == "table_row":
                    table_rows.append(items[i]["cells"])
                i += 1

            if table_rows:
                num_cols = max(len(r) for r in table_rows)
                tbl = doc.add_table(rows=len(table_rows), cols=num_cols)
                tbl.style = "Table Grid"
                for r_idx, row_cells in enumerate(table_rows):
                    is_header = (r_idx == 0)
                    for c_idx in range(num_cols):
                        val = row_cells[c_idx] if c_idx < len(row_cells) else ""
                        _set_cell(
                            tbl.cell(r_idx, c_idx),
                            val,
                            bold=is_header,
                            bg_hex="D9E1F2" if is_header else None,
                        )
                doc.add_paragraph("")  # 테이블 뒤 여백

        elif item["type"] == "table_sep":
            i += 1  # 단독 구분선은 무시

        else:  # text
            doc.add_paragraph(item["text"])
            i += 1

    doc.save(str(file_path))


def _add_pptx_run(p, item: dict) -> None:
    """pptx 단락(paragraph)에 item 내용을 run으로 추가한다."""
    from pptx.util import Pt

    if item["type"] == "bullet":
        # bullet은 level=0 으로 들여쓰기 유지 (레이아웃 기본 bullet 스타일 활용)
        p.level = 0
        run = p.add_run()
        run.text = item["text"]
    elif item["type"] == "heading" and item["level"] == 3:
        # ### 서브헤딩 → 볼드 + 약간 큰 폰트
        run = p.add_run()
        run.text = item["text"]
        run.font.bold = True
        run.font.size = Pt(14)
    else:
        run = p.add_run()
        run.text = item["text"]


def _save_pptx(content: str, file_path: Path) -> None:
    """
    pptx: # / ## 헤딩마다 슬라이드 1장 생성.

    슬라이드 구조:
      - 제목(title placeholder)   : # 또는 ## 헤딩 텍스트
      - 본문(body placeholder)    : bullet / text / ### 서브헤딩을 단락별로 분리 추가
      - 테이블(별도 shape)        : 마크다운 | 표 | 가 있으면 슬라이드 하단에 실제 테이블 추가

    헤딩 없으면 전체를 단일 슬라이드로 저장.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    TITLE_BODY_LAYOUT = 1  # 제목 + 내용 레이아웃

    items = _parse_md_lines(content)

    # ── 슬라이드 구분: level 1~2 헤딩마다 분리 ──────────────────────
    slides_data: list[dict] = []
    current: dict | None = None

    for item in items:
        if item["type"] == "heading" and item["level"] <= 2:
            if current is not None:
                slides_data.append(current)
            current = {"title": item["text"], "body_items": []}
        else:
            if current is None:
                current = {"title": "", "body_items": []}
            current["body_items"].append(item)

    if current is not None:
        slides_data.append(current)

    if not slides_data:
        slides_data = [{"title": "보고서", "body_items": [
            {"type": "text", "level": 0, "text": content, "cells": []}
        ]}]

    # ── 슬라이드 생성 ────────────────────────────────────────────────
    for sd in slides_data:
        body_items = sd["body_items"]
        has_table = any(it["type"] == "table_row" for it in body_items)
        # 테이블을 제외한 텍스트 항목
        text_items = [
            it for it in body_items
            if it["type"] not in ("table_row", "table_sep", "blank")
        ]

        slide = prs.slides.add_slide(prs.slide_layouts[TITLE_BODY_LAYOUT])
        slide.shapes.title.text = sd["title"]
        ph = slide.placeholders[1]  # body placeholder

        if has_table:
            # ── 테이블이 있는 슬라이드 ─────────────────────────────
            if text_items:
                # body placeholder를 상단 1/3 영역으로 축소
                ph.top    = Inches(1.55)
                ph.left   = Inches(0.5)
                ph.width  = Inches(9.0)
                ph.height = Inches(1.4)
                tf = ph.text_frame
                tf.clear()
                tf.word_wrap = True
                first = True
                for it in text_items:
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    _add_pptx_run(p, it)
            else:
                # 텍스트 없으면 body placeholder 최소화(사실상 숨김)
                ph.height = Pt(1)

            # 테이블 행 수집
            table_rows: list[list[str]] = []
            for it in body_items:
                if it["type"] == "table_row":
                    table_rows.append(it["cells"])

            if table_rows:
                num_rows = len(table_rows)
                num_cols = max(len(r) for r in table_rows)
                tbl_top  = Inches(3.1) if text_items else Inches(1.7)
                row_h    = Inches(0.42)

                tbl_shape = slide.shapes.add_table(
                    num_rows, num_cols,
                    Inches(0.5), tbl_top,
                    Inches(9.0), row_h * num_rows,
                )
                tbl = tbl_shape.table

                for r_idx, row_cells in enumerate(table_rows):
                    for c_idx in range(num_cols):
                        val  = row_cells[c_idx] if c_idx < len(row_cells) else ""
                        cell = tbl.cell(r_idx, c_idx)
                        cell.text_frame.clear()
                        cp  = cell.text_frame.paragraphs[0]
                        run = cp.add_run()
                        run.text = val
                        if r_idx == 0:          # 첫 행 = 헤더 → 볼드
                            run.font.bold = True

        else:
            # ── 텍스트만 있는 슬라이드 ────────────────────────────
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            first = True
            for it in body_items:
                if it["type"] in ("blank", "table_sep"):
                    continue
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                first = False
                _add_pptx_run(p, it)

    prs.save(str(file_path))


def _save_xlsx(content: str, file_path: Path) -> None:
    """
    xlsx: 마크다운 테이블(| col | col |)이 있으면 표로 저장,
    없으면 줄 단위로 A열에 순서대로 기록.
    헤딩은 굵게 표시.
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "보고서"

    lines = content.splitlines()

    # 마크다운 테이블 감지 (| 로 시작하는 줄이 2줄 이상 연속)
    table_lines = [l for l in lines if l.strip().startswith("|")]
    has_table = len(table_lines) >= 2

    if has_table:
        # 테이블 파싱 후 시트에 기록
        row_idx = 1
        for line in lines:
            stripped = line.strip()
            if stripped.startswith("|") and not set(stripped.replace("|", "").replace("-", "").replace(" ", "")) == set():
                # 구분선(|---|---|) 제외
                if all(c in "-| :" for c in stripped):
                    continue
                cells = [c.strip() for c in stripped.strip("|").split("|")]
                for col_idx, cell_val in enumerate(cells, start=1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=cell_val)
                    if row_idx == 1:
                        cell.font = Font(bold=True)
                        cell.fill = PatternFill("solid", fgColor="D9E1F2")
                    cell.alignment = Alignment(wrap_text=True)
                row_idx += 1
            elif stripped and not stripped.startswith("|"):
                # 테이블 밖 텍스트도 기록
                ws.cell(row=row_idx, column=1, value=stripped)
                row_idx += 1

        # 열 너비 자동 조정
        for col in ws.columns:
            max_len = max((len(str(c.value or "")) for c in col), default=10)
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_len + 4, 60)
    else:
        # 줄 단위 기록
        for row_idx, item in enumerate(_parse_md_lines(content), start=1):
            cell = ws.cell(row=row_idx, column=1, value=item["text"])
            if item["type"] == "heading":
                size = {1: 16, 2: 14, 3: 12}.get(item["level"], 12)
                cell.font = Font(bold=True, size=size)
                if item["level"] == 1:
                    cell.fill = PatternFill("solid", fgColor="D9E1F2")
            elif item["type"] == "bullet":
                cell.value = "• " + item["text"]

        ws.column_dimensions["A"].width = 80

    wb.save(str(file_path))


# ══════════════════════════════════════════════════════════════════
#  보고서 저장 도구
# ══════════════════════════════════════════════════════════════════

@tool
def save_report(
    content: str,
    filename: str = "report",
    format: Literal["md", "txt", "docx", "pptx", "xlsx"] = "md",
    output_dir: str = "",
    add_timestamp: bool = True,
) -> str:
    """
    내용을 보고서 파일로 저장한다.
    에이전트가 생성한 요약·분석·번역 결과를 다양한 형식으로 저장할 때 사용.

    Args:
        content      : 저장할 내용 (마크다운 형식 권장)
        filename     : 파일명 (확장자 제외, 기본: report)
        format       : 파일 형식
                       md   — 마크다운 텍스트
                       txt  — 일반 텍스트
                       docx — Word 문서 (헤딩·글머리기호 자동 스타일 적용)
                       pptx — PowerPoint (## 헤딩마다 슬라이드 1장)
                       xlsx — Excel (표 구조 자동 감지, 없으면 줄 단위 기록)
        output_dir   : 저장 디렉터리 경로 (기본: AGENT_DATA_DIR/reports)
        add_timestamp: True면 파일명에 날짜시간 추가 (기본: True)
    """
    from datetime import datetime

    output_path = Path(output_dir) if output_dir else _DEFAULT_OUTPUT_DIR
    output_path.mkdir(parents=True, exist_ok=True)

    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    final_name = f"{filename}_{ts}.{format}" if add_timestamp else f"{filename}.{format}"
    file_path = output_path / final_name

    handler_map = {
        "txt":  _save_txt,
        "md":   _save_txt,
        "docx": _save_docx,
        "pptx": _save_pptx,
        "xlsx": _save_xlsx,
    }

    try:
        handler_map[format](content, file_path)
        abs_path = file_path.resolve()
        size_kb = abs_path.stat().st_size / 1024
        print(f"[FileTools] 보고서 저장 완료: {abs_path}")
        return (
            f"[save_report] 저장 완료\n"
            f"  경로: {abs_path}\n"
            f"  크기: {size_kb:.1f} KB\n"
            f"  형식: {format.upper()}"
        )
    except Exception as exc:
        return f"[Error] 파일 저장 실패 ({format}): {exc}"


# ══════════════════════════════════════════════════════════════════
#  도구 목록
# ══════════════════════════════════════════════════════════════════

file_tools = [
    load_file_content,
    summarize_file,
    translate_file,
    extract_key_points,
    analyze_file,
    save_report,
]
